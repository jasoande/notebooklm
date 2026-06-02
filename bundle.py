import os
import csv
import time
import shutil
import subprocess
from pptx import Presentation
from concurrent.futures import ProcessPoolExecutor
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

# =====================================================================
# CONFIGURATION
# =====================================================================
SOURCE_DIR = "."
WORKSPACE_DIR = os.path.abspath("./pdf_staging_pool")
MASTER_PDF = "master.pdf"
MASTER_CSV = "master.csv"

# Absolute exclusions to prevent infinite loops
EXCLUDED_FILES = { "bundle.py", "pull.py", MASTER_PDF, MASTER_CSV }

# Detect macOS LibreOffice bin setups cleanly
LIBRE_EXEC = "libreoffice"
if os.path.exists(os.path.expanduser("~/bin/libreoffice")):
    LIBRE_EXEC = os.path.expanduser("~/bin/libreoffice")
elif os.path.exists("/usr/local/bin/libreoffice"):
    LIBRE_EXEC = "/usr/local/bin/libreoffice"

# =====================================================================
# CORE CONVERSION ENGINES (Worker Safe)
# =====================================================================
def convert_pptx_to_pdf(pptx_path):
    """Extracts slide frameworks and compiles them to a document tree."""
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_pdf = os.path.join(WORKSPACE_DIR, f"2_pptx_{base_name}.pdf")
    
    doc = SimpleDocTemplate(output_pdf, pagesize=letter,
                            rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    story = []
    styles = getSampleStyleSheet()
    
    prs = Presentation(pptx_path)
    story.append(Paragraph(f"Presentation Deck: {base_name}", styles['Heading1']))
    story.append(Spacer(1, 20))
    
    for i, slide in enumerate(prs.slides):
        story.append(Paragraph(f"Slide {i+1}", styles['Heading3']))
        story.append(Spacer(1, 5))
        
        slide_text_blocks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text_blocks.append(shape.text.strip())
                
        if slide_text_blocks:
            combined_text = "<br/><br/>".join(slide_text_blocks)
            story.append(Paragraph(combined_text, styles['Normal']))
        else:
            story.append(Paragraph("[Visual Media Only Slide]", styles['Italic']))
        story.append(Spacer(1, 15))
        
    doc.build(story)
    print(f" -> Compiled PPTX to PDF: {base_name}")
    return output_pdf


def convert_office_via_libre(task_info):
    """Converts a single Office document using robust retry states to prevent lock crashes."""
    file_path, prefix, worker_id = task_info
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    
    # Isolate user profile variables completely away from file names
    unique_env_dir = os.path.join(WORKSPACE_DIR, f"env_worker_{worker_id}")
    unique_env_url = f"file://{unique_env_dir}"
    abs_file_path = os.path.abspath(file_path)
    
    cmd = [
        LIBRE_EXEC,
        f"-env:UserInstallation={unique_env_url}",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", WORKSPACE_DIR,
        abs_file_path
    ]
    
    # IMPLEMENTATION: Stateful retry logic to catch rapid IPC crashes
    max_retries = 3
    backoff_delay = 0.2  # Base sleep delay in seconds
    
    for attempt in range(1, max_retries + 1):
        try:
            # Clean up the specific workspace directory if a previous attempt corrupted it
            if os.path.exists(unique_env_dir):
                shutil.rmtree(unique_env_dir)
                
            subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            
            # If execution passes cleanly, break the retry cycle immediately
            break
        except subprocess.CalledProcessError as err:
            if attempt == max_retries:
                print(f" ! Permanent Failure converting {base_name} after {max_retries} attempts.")
                raise err
            print(f" [Collision Warning] Quick-fire crash detected for {base_name}. Retrying in {backoff_delay}s... (Attempt {attempt}/{max_retries})")
            time.sleep(backoff_delay)
            backoff_delay *= 2  # Double the delay time for the next attempt (exponential backoff)

    default_output = os.path.join(WORKSPACE_DIR, f"{base_name}.pdf")
    ordered_output = os.path.join(WORKSPACE_DIR, f"{prefix}_{base_name}.pdf")
    
    if os.path.exists(default_output):
        os.rename(default_output, ordered_output)
        
    if os.path.exists(unique_env_dir):
        shutil.rmtree(unique_env_dir)
        
    print(f" -> Compiled Office File ({prefix.upper()}) to PDF: {base_name}")
    return ordered_output


# =====================================================================
# SYSTEM MAIN ENGINE
# =====================================================================
def main():
    if subprocess.run("command -v pdfunite", shell=True, capture_output=True).returncode != 0:
        print("Error: 'pdfunite' is required but missing. Install it via 'brew install poppler'.")
        return
    if subprocess.run(f"command -v {LIBRE_EXEC}", shell=True, capture_output=True).returncode != 0:
        print(f"Error: '{LIBRE_EXEC}' is required for Word/Excel parsing but missing.")
        return

    print("Step 1: Scanning local filesystem architecture...")
    
    csv_files = sorted([f for f in os.listdir(SOURCE_DIR) if f.endswith('.csv') and f not in EXCLUDED_FILES])
    pptx_files = sorted([f for f in os.listdir(SOURCE_DIR) if f.endswith('.pptx') and f not in EXCLUDED_FILES])
    docx_files = sorted([f for f in os.listdir(SOURCE_DIR) if f.endswith('.docx') and f not in EXCLUDED_FILES])
    xlsx_files = sorted([f for f in os.listdir(SOURCE_DIR) if f.endswith('.xlsx') and f not in EXCLUDED_FILES])
    existing_pdfs = sorted([f for f in os.listdir(SOURCE_DIR) if f.endswith('.pdf') and f not in EXCLUDED_FILES])

    # -----------------------------------------------------------------
    # TRACK 1: CONSOLIDATING CSV DATA ROWS INTO master.csv
    # -----------------------------------------------------------------
    if csv_files:
        print(f"==> Track A: Merging {len(csv_files)} CSV data arrays into {MASTER_CSV}...")
        with open(MASTER_CSV, mode='w', encoding='utf-8', newline='') as master_f:
            writer = csv.writer(master_f)
            for idx, csv_file in enumerate(csv_files):
                file_path = os.path.join(SOURCE_DIR, csv_file)
                with open(file_path, mode='r', encoding='utf-8', errors='replace', newline='') as input_f:
                    reader = csv.reader(input_f)
                    for row_idx, row in enumerate(reader):
                        if row_idx == 0 and idx > 0:
                            continue  
                        writer.writerow(row)
        print(f" -> Success: {MASTER_CSV} compiled natively.")
    else:
        print(" ! Track A: No separate CSV file targets found.")

    # -----------------------------------------------------------------
    # TRACK 2: CONSOLIDATING PDF/PPTX/DOCX/XLSX INTO master.pdf
    # -----------------------------------------------------------------
    if pptx_files or docx_files or xlsx_files or existing_pdfs:
        print(f"==> Track B: Setting up parallel PDF workspace allocation...")
        
        if os.path.exists(WORKSPACE_DIR):
            shutil.rmtree(WORKSPACE_DIR)
        os.makedirs(WORKSPACE_DIR, exist_ok=True)

        for pdf in existing_pdfs:
            shutil.copy(os.path.join(SOURCE_DIR, pdf), os.path.join(WORKSPACE_DIR, f"1_std_{pdf}"))

        print(" -> Spawning parallel compilation workers with lock safeguards...")
        
        libre_tasks = []
        worker_id = 0
        for f in docx_files:
            libre_tasks.append((os.path.join(SOURCE_DIR, f), "3_docx", worker_id))
            worker_id += 1
        for f in xlsx_files:
            libre_tasks.append((os.path.join(SOURCE_DIR, f), "4_xlsx", worker_id))
            worker_id += 1

        # MODIFIED: Throttling max_workers to 2 prevents LibreOffice instances 
        # from stepping on each other's memory frames.
        with ProcessPoolExecutor(max_workers=2) as executor:
            futures = []
            
            for f in pptx_files:
                futures.append(executor.submit(convert_pptx_to_pdf, os.path.join(SOURCE_DIR, f)))
            for task in libre_tasks:
                futures.append(executor.submit(convert_office_via_libre, task))
                
            for future in futures:
                future.result()

        final_pdf_pool = [os.path.join(WORKSPACE_DIR, f) for f in os.listdir(WORKSPACE_DIR) if f.endswith('.pdf')]
        final_pdf_pool.sort()

        if final_pdf_pool:
            print(f" -> Executing native pdfunite consolidation into {MASTER_PDF}...")
            try:
                if os.path.exists(MASTER_PDF):
                    os.remove(MASTER_PDF)
                    
                cmd = ["pdfunite"] + final_pdf_pool + [MASTER_PDF]
                subprocess.run(cmd, check=True)
                
                shutil.rmtree(WORKSPACE_DIR)
                print(f" -> Success: {MASTER_PDF} compiled natively.")
                
            except subprocess.CalledProcessError as err:
                print(f"Compilation Error during pdfunite phase: {err}")
    else:
        print(" ! Track B: No PDF, PPTX, DOCX, or XLSX asset targets found.")

    print("\n==> Success Pipeline Completed! Master consolidation vectors are finalized.")

if __name__ == '__main__':
    main()
